import sys, os, json, csv, zipfile, traceback, hashlib, mimetypes, difflib, tempfile, re, io, contextlib
from pathlib import Path
from datetime import datetime


def out(obj):
    sys.stdout.write(json.dumps(obj, ensure_ascii=False, default=str))
    sys.stdout.flush()


def sha256_file(path: Path):
    h = hashlib.sha256()
    with path.open('rb') as f:
        for chunk in iter(lambda: f.read(1024 * 1024), b''):
            h.update(chunk)
    return h.hexdigest()


def base_meta(path: Path):
    stat = path.stat()
    mime, _ = mimetypes.guess_type(path.name)
    return {
        "success": True,
        "action": "inspect",
        "path": str(path),
        "name": path.name,
        "extension": path.suffix.lower(),
        "bytes": stat.st_size,
        "modified_at": datetime.fromtimestamp(stat.st_mtime).isoformat(timespec='seconds'),
        "mime": mime or "application/octet-stream",
        "sha256": sha256_file(path),
    }


def safe_text(v):
    if v is None:
        return ""
    return str(v)


def inspect_pdf(path: Path, base):
    import fitz
    doc = fitz.open(str(path))
    pages = []
    total_chars = 0
    total_images = 0
    scanned_candidates = []
    table_count = 0
    for i, page in enumerate(doc):
        if i >= 120:
            break
        text = page.get_text('text') or ''
        images = len(page.get_images(full=True))
        links = len(page.get_links())
        total_chars += len(text)
        total_images += images
        tables = []
        try:
            with contextlib.redirect_stdout(io.StringIO()):
                finder = page.find_tables()
            for table in getattr(finder, 'tables', [])[:10]:
                extracted = table.extract()
                tables.append(extracted[:40])
            table_count += len(tables)
        except Exception:
            pass
        if len(text.strip()) < 40 and images > 0:
            scanned_candidates.append(i + 1)
        pages.append({
            "page": i + 1,
            "text": text[:24000],
            "chars": len(text),
            "words": len(text.split()),
            "images": images,
            "links": links,
            "tables": tables,
            "width": round(page.rect.width, 2),
            "height": round(page.rect.height, 2),
        })
    metadata = {}
    try:
        metadata = {k: v for k, v in (doc.metadata or {}).items() if v}
    except Exception:
        pass
    base.update({
        "type": "pdf",
        "page_count": len(doc),
        "pages": pages,
        "metadata": metadata,
        "text_chars": total_chars,
        "image_count": total_images,
        "table_count": table_count,
        "scanned_candidate_pages": scanned_candidates[:100],
        "needs_ocr": bool(scanned_candidates),
        "truncated": len(doc) > 120,
    })
    doc.close()
    return base


def inspect_docx(path: Path, base):
    from docx import Document
    doc = Document(str(path))
    paragraphs = []
    headings = []
    for idx, p in enumerate(doc.paragraphs):
        text = p.text.strip()
        if not text:
            continue
        item = {"index": idx, "text": text[:12000], "style": p.style.name if p.style else ""}
        paragraphs.append(item)
        if item["style"].lower().startswith('heading'):
            headings.append(item)
    tables = []
    for ti, table in enumerate(doc.tables[:60]):
        rows = [[cell.text for cell in row.cells[:60]] for row in table.rows[:250]]
        tables.append({"table": ti + 1, "rows": rows, "row_count": len(table.rows), "column_count": len(table.columns)})
    props = doc.core_properties
    base.update({
        "type": "word",
        "paragraphs": paragraphs[:3500],
        "headings": headings[:500],
        "tables": tables,
        "paragraph_count": len(doc.paragraphs),
        "table_count": len(doc.tables),
        "image_count": len(doc.inline_shapes),
        "metadata": {
            "title": props.title or "", "subject": props.subject or "", "author": props.author or "",
            "keywords": props.keywords or "", "created": props.created, "modified": props.modified,
        },
        "truncated": len(paragraphs) > 3500 or len(doc.tables) > 60,
    })
    return base


def inspect_xlsx(path: Path, base):
    import openpyxl
    keep_vba = path.suffix.lower() == '.xlsm'
    wb = openpyxl.load_workbook(str(path), data_only=False, read_only=False, keep_vba=keep_vba)
    sheets = []
    formula_count = 0
    nonempty_count = 0
    for ws in wb.worksheets[:40]:
        rows = []
        local_formula_count = 0
        for ri, row in enumerate(ws.iter_rows(values_only=False)):
            if ri >= 500:
                break
            vals = []
            for cell in row[:80]:
                val = cell.value
                if val not in (None, ''):
                    nonempty_count += 1
                if isinstance(val, str) and val.startswith('='):
                    formula_count += 1
                    local_formula_count += 1
                vals.append(val)
            rows.append(vals)
        sheets.append({
            "name": ws.title,
            "state": ws.sheet_state,
            "max_row": ws.max_row,
            "max_column": ws.max_column,
            "formula_count": local_formula_count,
            "merged_ranges": [str(x) for x in list(ws.merged_cells.ranges)[:100]],
            "freeze_panes": str(ws.freeze_panes or ''),
            "rows": rows,
            "truncated": ws.max_row > 500 or ws.max_column > 80,
        })
    base.update({
        "type": "excel",
        "sheet_count": len(wb.sheetnames),
        "sheet_names": wb.sheetnames,
        "sheets": sheets,
        "formula_count": formula_count,
        "nonempty_cell_count_sampled": nonempty_count,
        "defined_names": [str(x) for x in list(wb.defined_names.values())[:100]],
        "truncated": len(wb.sheetnames) > 40,
    })
    wb.close()
    return base


def inspect_csv(path: Path, base):
    rows = []
    chosen_encoding = None
    last_error = None
    for enc in ('utf-8-sig', 'utf-8', 'cp1256', 'windows-1252', 'latin-1'):
        try:
            with open(path, 'r', encoding=enc, newline='') as f:
                reader = csv.reader(f)
                for i, row in enumerate(reader):
                    if i >= 1000:
                        break
                    rows.append(row[:100])
            chosen_encoding = enc
            break
        except Exception as e:
            rows = []
            last_error = e
    if chosen_encoding is None:
        raise RuntimeError(f"تعذر قراءة CSV: {last_error}")
    widths = [len(r) for r in rows] or [0]
    base.update({
        "type": "csv", "rows": rows, "row_count_sampled": len(rows),
        "max_columns_sampled": max(widths), "encoding": chosen_encoding, "truncated": len(rows) >= 1000,
    })
    return base


def inspect_pptx(path: Path, base):
    from pptx import Presentation
    prs = Presentation(str(path))
    slides = []
    total_pictures = 0
    total_tables = 0
    for i, slide in enumerate(prs.slides):
        if i >= 150:
            break
        texts = []
        pictures = 0
        tables = []
        title = ""
        try:
            if slide.shapes.title:
                title = slide.shapes.title.text or ""
        except Exception:
            pass
        for shape in slide.shapes:
            if getattr(shape, 'has_text_frame', False) and shape.text:
                texts.append(shape.text)
            if getattr(shape, 'shape_type', None) == 13:  # MSO_SHAPE_TYPE.PICTURE
                pictures += 1
            if getattr(shape, 'has_table', False):
                t = [[cell.text for cell in row.cells] for row in list(shape.table.rows)[:80]]
                tables.append(t)
        total_pictures += pictures
        total_tables += len(tables)
        slides.append({"slide": i + 1, "title": title, "text": "\n".join(texts)[:24000], "pictures": pictures, "tables": tables[:15]})
    base.update({
        "type": "powerpoint", "slide_count": len(prs.slides), "slides": slides,
        "picture_count": total_pictures, "table_count": total_tables, "truncated": len(prs.slides) > 150,
    })
    return base


def inspect_zip(path: Path, base):
    with zipfile.ZipFile(path, 'r') as z:
        infos = z.infolist()
        names = [i.filename for i in infos]
        ext_counts = {}
        total_uncompressed = 0
        for info in infos:
            total_uncompressed += info.file_size
            ext = Path(info.filename).suffix.lower() or '[none]'
            ext_counts[ext] = ext_counts.get(ext, 0) + 1
        markers = [n for n in names if Path(n).name.lower() in {
            'package.json','requirements.txt','pyproject.toml','cargo.toml','go.mod','pom.xml','composer.json','readme.md','dockerfile'
        }]
    base.update({
        "type": "zip", "entries": names[:3000], "entry_count": len(names),
        "extension_counts": dict(sorted(ext_counts.items(), key=lambda x: x[1], reverse=True)[:60]),
        "project_markers": markers[:100], "uncompressed_bytes": total_uncompressed,
        "truncated": len(names) > 3000,
    })
    return base


def inspect_image(path: Path, base):
    from PIL import Image, ExifTags
    with Image.open(path) as im:
        exif = {}
        try:
            raw = im.getexif()
            for k, v in list(raw.items())[:80]:
                key = ExifTags.TAGS.get(k, str(k))
                exif[key] = safe_text(v)[:500]
        except Exception:
            pass
        base.update({
            "type": "image", "width": im.width, "height": im.height, "mode": im.mode,
            "format": im.format, "frames": getattr(im, 'n_frames', 1), "exif": exif,
        })
    return base


def inspect_file(p: str):
    path = Path(p).resolve()
    if not path.exists():
        raise FileNotFoundError(str(path))
    if not path.is_file():
        raise RuntimeError('المسار ليس ملفًا.')
    ext = path.suffix.lower()
    base = base_meta(path)

    if ext == '.pdf':
        return inspect_pdf(path, base)
    if ext == '.docx':
        return inspect_docx(path, base)
    if ext in ('.xlsx', '.xlsm'):
        return inspect_xlsx(path, base)
    if ext == '.csv':
        return inspect_csv(path, base)
    if ext == '.pptx':
        return inspect_pptx(path, base)
    if ext == '.zip':
        return inspect_zip(path, base)
    if ext in ('.png','.jpg','.jpeg','.webp','.bmp','.gif','.tif','.tiff'):
        return inspect_image(path, base)

    # Legacy formats that need Microsoft Office/LibreOffice conversion before deep parsing.
    if ext in ('.doc', '.xls', '.ppt', '.rtf', '.odt', '.ods', '.odp'):
        base.update({"type": "legacy_office", "note": "هذا تنسيق قديم/خارجي. التحليل العميق يحتاج تحويله إلى DOCX/XLSX/PPTX أو PDF أولًا."})
        return base

    try:
        text = path.read_text(encoding='utf-8', errors='replace')
        base.update({"type": "text", "content": text[:400000], "chars": len(text), "truncated": len(text) > 400000})
    except Exception:
        base.update({"type": "binary", "note": "لا يوجد Parser متخصص لهذا الامتداد حتى الآن."})
    return base


def inspection_text(obj):
    t = obj.get('type')
    if t == 'pdf':
        return '\n'.join(f"[Page {p.get('page')}]\n{p.get('text','')}" for p in obj.get('pages', []))
    if t == 'word':
        pars = '\n'.join(p.get('text','') for p in obj.get('paragraphs', []))
        tabs = '\n'.join('\t'.join(safe_text(c) for c in row) for tb in obj.get('tables', []) for row in tb.get('rows', []))
        return pars + '\n' + tabs
    if t == 'excel':
        parts = []
        for sh in obj.get('sheets', []):
            parts.append(f"[Sheet {sh.get('name')}]\n" + '\n'.join('\t'.join(safe_text(c) for c in row) for row in sh.get('rows', [])))
        return '\n'.join(parts)
    if t == 'csv':
        return '\n'.join('\t'.join(safe_text(c) for c in row) for row in obj.get('rows', []))
    if t == 'powerpoint':
        return '\n'.join(f"[Slide {s.get('slide')}] {s.get('title','')}\n{s.get('text','')}" for s in obj.get('slides', []))
    if t == 'text':
        return obj.get('content', '')
    if t == 'zip':
        return '\n'.join(obj.get('entries', []))
    return json.dumps(obj, ensure_ascii=False, default=str)


def line_key(line):
    return re.sub(r'\s+', ' ', line.strip())


def compare_files(payload):
    paths = [Path(p).resolve() for p in payload.get('paths', [])]
    if len(paths) < 2:
        raise RuntimeError('compare يحتاج ملفين على الأقل.')
    inspections = [inspect_file(str(p)) for p in paths[:5]]
    base_text = inspection_text(inspections[0])
    base_lines = [line_key(x) for x in base_text.splitlines() if line_key(x)]
    comparisons = []
    for other in inspections[1:]:
        other_text = inspection_text(other)
        other_lines = [line_key(x) for x in other_text.splitlines() if line_key(x)]
        ratio = difflib.SequenceMatcher(None, base_text[:500000], other_text[:500000]).ratio()
        a_set, b_set = set(base_lines), set(other_lines)
        only_a = [x for x in base_lines if x not in b_set][:120]
        only_b = [x for x in other_lines if x not in a_set][:120]
        comparisons.append({
            "against": other.get('path'), "similarity": round(ratio, 4),
            "only_in_base": only_a, "only_in_other": only_b,
            "base_line_count": len(base_lines), "other_line_count": len(other_lines),
        })
    return {"success": True, "action": "compare", "base": inspections[0], "files": inspections, "comparisons": comparisons}


def search_files(payload):
    paths = [Path(p).resolve() for p in payload.get('paths', [])][:30]
    query = str(payload.get('query') or '').strip()
    if not query:
        return {"success": True, "action": "search", "query": query, "results": []}
    q = query.casefold()
    max_results = int(payload.get('max_results') or 40)
    results = []
    for path in paths:
        try:
            obj = inspect_file(str(path))
            text = inspection_text(obj)
            low = text.casefold()
            start = 0
            hits = 0
            while len(results) < max_results:
                idx = low.find(q, start)
                if idx < 0:
                    break
                snippet = text[max(0, idx - 180): min(len(text), idx + len(query) + 260)].replace('\x00', '')
                results.append({"path": str(path), "name": path.name, "type": obj.get('type'), "snippet": snippet})
                start = idx + max(1, len(query))
                hits += 1
                if hits >= 8:
                    break
        except Exception as e:
            results.append({"path": str(path), "name": path.name, "error": str(e)})
        if len(results) >= max_results:
            break
    return {"success": True, "action": "search", "query": query, "results": results[:max_results]}


def render_pdf_pages(payload):
    import fitz
    path = Path(payload['path']).resolve()
    if path.suffix.lower() != '.pdf':
        raise RuntimeError('render_pdf_pages يدعم PDF فقط.')
    pages = payload.get('pages') or [1]
    pages = [int(x) for x in pages][:10]
    scale = float(payload.get('scale') or 1.7)
    out_dir = Path(payload.get('output_dir') or tempfile.mkdtemp(prefix='abdx_pdf_')).resolve()
    out_dir.mkdir(parents=True, exist_ok=True)
    doc = fitz.open(str(path))
    outputs = []
    for n in pages:
        if n < 1 or n > len(doc):
            continue
        page = doc[n - 1]
        pix = page.get_pixmap(matrix=fitz.Matrix(scale, scale), alpha=False)
        out_path = out_dir / f"page-{n}.png"
        pix.save(str(out_path))
        outputs.append({"page": n, "path": str(out_path), "width": pix.width, "height": pix.height})
    doc.close()
    return {"success": True, "action": "render_pdf_pages", "source": str(path), "pages": outputs}


def _hex_color(value, default='1F4E78'):
    text = str(value or default).strip().replace('#', '')
    return text.upper() if re.fullmatch(r'[0-9A-Fa-f]{6}', text) else default


def _set_cell_shading(cell, color):
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = tc_pr.find(qn('w:shd'))
    if shd is None:
        shd = OxmlElement('w:shd')
        tc_pr.append(shd)
    shd.set(qn('w:fill'), _hex_color(color))


def _set_paragraph_rtl(paragraph, rtl=True):
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
    p_pr = paragraph._p.get_or_add_pPr()
    bidi = p_pr.find(qn('w:bidi'))
    if rtl and bidi is None:
        bidi = OxmlElement('w:bidi')
        bidi.set(qn('w:val'), '1')
        p_pr.append(bidi)
    elif not rtl and bidi is not None:
        p_pr.remove(bidi)


def _add_word_page_number(paragraph):
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
    run = paragraph.add_run()
    fld_char1 = OxmlElement('w:fldChar'); fld_char1.set(qn('w:fldCharType'), 'begin')
    instr = OxmlElement('w:instrText'); instr.set(qn('xml:space'), 'preserve'); instr.text = ' PAGE '
    fld_char2 = OxmlElement('w:fldChar'); fld_char2.set(qn('w:fldCharType'), 'end')
    run._r.extend([fld_char1, instr, fld_char2])


def _word_apply_branding(doc, payload):
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.section import WD_SECTION
    brand = payload.get('branding') or {}
    primary = _hex_color(brand.get('primary_color'), '1F4E78')
    font_name = brand.get('font') or 'Arial'
    rtl = bool(payload.get('rtl', brand.get('rtl', False)))
    for sec in doc.sections:
        margins = payload.get('margins') or {}
        sec.top_margin = Inches(float(margins.get('top', 0.7)))
        sec.bottom_margin = Inches(float(margins.get('bottom', 0.7)))
        sec.left_margin = Inches(float(margins.get('left', 0.7)))
        sec.right_margin = Inches(float(margins.get('right', 0.7)))
        header_text = brand.get('header_text') or ''
        if header_text:
            hp = sec.header.paragraphs[0]
            hp.text = header_text
            hp.alignment = WD_ALIGN_PARAGRAPH.RIGHT if rtl else WD_ALIGN_PARAGRAPH.LEFT
            _set_paragraph_rtl(hp, rtl)
        footer_text = brand.get('footer_text') or ''
        fp = sec.footer.paragraphs[0]
        if footer_text:
            fp.add_run(footer_text + '  |  ')
        if payload.get('page_numbers', True):
            _add_word_page_number(fp)
        fp.alignment = WD_ALIGN_PARAGRAPH.CENTER
    styles = doc.styles
    for style_name in ('Normal', 'Title', 'Heading 1', 'Heading 2', 'Heading 3'):
        if style_name in styles:
            st = styles[style_name]
            st.font.name = font_name
            if style_name.startswith('Heading'):
                st.font.color.rgb = RGBColor.from_string(primary)
    if 'Normal' in styles:
        styles['Normal'].font.size = Pt(float(brand.get('body_size', 11)))
    return primary, font_name, rtl


def _word_add_table(doc, table_spec, primary, rtl=False):
    from docx.shared import Pt
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    headers = table_spec.get('headers') or []
    rows = table_spec.get('rows') or []
    cols = max(len(headers), max((len(r) for r in rows), default=0), 1)
    table = doc.add_table(rows=1 if headers else 0, cols=cols)
    table.style = table_spec.get('style') or 'Table Grid'
    if headers:
        for i, value in enumerate(headers[:cols]):
            c = table.rows[0].cells[i]
            c.text = safe_text(value)
            _set_cell_shading(c, table_spec.get('header_color') or primary)
            for run in c.paragraphs[0].runs:
                run.font.bold = True
                try:
                    from docx.shared import RGBColor
                    run.font.color.rgb = RGBColor(255,255,255)
                except Exception: pass
            c.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.RIGHT if rtl else WD_ALIGN_PARAGRAPH.LEFT
            _set_paragraph_rtl(c.paragraphs[0], rtl)
    for row in rows:
        cells = table.add_row().cells
        for i in range(cols):
            cells[i].text = safe_text(row[i] if i < len(row) else '')
            cells[i].paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.RIGHT if rtl else WD_ALIGN_PARAGRAPH.LEFT
            _set_paragraph_rtl(cells[i].paragraphs[0], rtl)
    return table


def create_word(payload):
    from docx import Document
    from docx.shared import Inches, Pt
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    path = Path(payload['path']); path.parent.mkdir(parents=True, exist_ok=True)
    template = payload.get('template_path')
    doc = Document(str(template)) if template and Path(template).exists() else Document()
    primary, font_name, rtl = _word_apply_branding(doc, payload)

    logo = (payload.get('branding') or {}).get('logo_path')
    if logo and Path(logo).exists():
        p_logo = doc.add_paragraph(); p_logo.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p_logo.add_run().add_picture(str(logo), width=Inches(float((payload.get('branding') or {}).get('logo_width', 1.4))))

    title = doc.add_paragraph(); title.style = doc.styles['Title'] if 'Title' in doc.styles else None
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER; _set_paragraph_rtl(title, rtl)
    run = title.add_run(payload.get('title','Report')); run.bold = True; run.font.size = Pt(22)
    subtitle = payload.get('subtitle')
    if subtitle:
        p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER; _set_paragraph_rtl(p, rtl)
        p.add_run(subtitle).italic = True
    if payload.get('cover_page', False):
        doc.add_page_break()

    for item in payload.get('sections', []):
        heading = item.get('heading','')
        if heading:
            hp = doc.add_heading(heading, level=int(item.get('level', 1) or 1)); _set_paragraph_rtl(hp, rtl)
        for para in str(item.get('body','')).split('\n'):
            if para.strip():
                p = doc.add_paragraph(para); p.paragraph_format.space_after = Pt(6)
                p.alignment = WD_ALIGN_PARAGRAPH.RIGHT if rtl else WD_ALIGN_PARAGRAPH.LEFT; _set_paragraph_rtl(p, rtl)
        for bullet in item.get('bullets') or []:
            p = doc.add_paragraph(style='List Bullet'); p.add_run(safe_text(bullet)); _set_paragraph_rtl(p, rtl)
        for table_spec in item.get('tables') or []:
            _word_add_table(doc, table_spec, primary, rtl)
        for image in item.get('images') or []:
            image_path = image.get('path') if isinstance(image, dict) else image
            if image_path and Path(image_path).exists():
                p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                width = float(image.get('width_inches', 5.5)) if isinstance(image, dict) else 5.5
                p.add_run().add_picture(str(image_path), width=Inches(width))
                if isinstance(image, dict) and image.get('caption'):
                    cap = doc.add_paragraph(image['caption']); cap.alignment = WD_ALIGN_PARAGRAPH.CENTER; _set_paragraph_rtl(cap, rtl)
    for table_spec in payload.get('tables') or []:
        _word_add_table(doc, table_spec, primary, rtl)
    doc.save(str(path))
    return {"success":True,"action":"create_word","path":str(path),"sections":len(payload.get('sections',[])),"tables":len(doc.tables),"rtl":rtl,"template_used":bool(template and Path(template).exists())}


def edit_word(payload):
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Pt
    path = Path(payload['path']).resolve()
    if not path.exists(): raise FileNotFoundError(str(path))
    out_path = Path(payload.get('output_path') or path).resolve(); out_path.parent.mkdir(parents=True, exist_ok=True)
    doc = Document(str(path)); primary, font_name, rtl = _word_apply_branding(doc, payload)
    replacements = payload.get('replacements') or []
    replacement_count = 0
    for rep in replacements:
        old, new = safe_text(rep.get('find')), safe_text(rep.get('replace'))
        if not old: continue
        for p in doc.paragraphs:
            if old in p.text:
                for run in p.runs:
                    if old in run.text:
                        run.text = run.text.replace(old, new); replacement_count += 1
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for p in cell.paragraphs:
                        if old in p.text:
                            for run in p.runs:
                                if old in run.text:
                                    run.text = run.text.replace(old, new); replacement_count += 1
    for sec in payload.get('append_sections') or []:
        if sec.get('heading'):
            hp=doc.add_heading(sec.get('heading'), level=int(sec.get('level',1) or 1)); _set_paragraph_rtl(hp, rtl)
        for para in str(sec.get('body','')).split('\n'):
            if para.strip():
                pp=doc.add_paragraph(para); pp.paragraph_format.space_after=Pt(6); _set_paragraph_rtl(pp, rtl)
        for tb in sec.get('tables') or []: _word_add_table(doc, tb, primary, rtl)
    for tb in payload.get('append_tables') or []: _word_add_table(doc, tb, primary, rtl)
    doc.save(str(out_path))
    return {"success":True,"action":"edit_word","path":str(out_path),"replacements":replacement_count,"tables":len(doc.tables),"rtl":rtl}


def _excel_style_header(ws, row_num, start_col, end_col, primary):
    from openpyxl.styles import Font, PatternFill, Alignment
    fill = PatternFill('solid', fgColor=primary)
    for col in range(start_col, end_col + 1):
        c = ws.cell(row_num, col); c.font = Font(bold=True, color='FFFFFF'); c.fill = fill; c.alignment = Alignment(horizontal='center', vertical='center')


def _excel_autofit(ws, max_rows=300, max_cols=80):
    from openpyxl.utils import get_column_letter
    for col_idx in range(1, min(ws.max_column, max_cols) + 1):
        max_len = 0
        for r in range(1, min(ws.max_row, max_rows) + 1):
            v = ws.cell(r, col_idx).value
            max_len = max(max_len, len(str(v)) if v is not None else 0)
        ws.column_dimensions[get_column_letter(col_idx)].width = min(max(max_len + 2, 10), 45)


def _excel_add_chart(ws, spec):
    from openpyxl.chart import BarChart, LineChart, PieChart, Reference
    kind = str(spec.get('type') or 'bar').lower()
    chart = PieChart() if kind == 'pie' else (LineChart() if kind == 'line' else BarChart())
    chart.title = spec.get('title') or 'Chart'
    min_row = int(spec.get('min_row', 1)); max_row = int(spec.get('max_row', ws.max_row))
    min_col = int(spec.get('min_col', 1)); max_col = int(spec.get('max_col', min(ws.max_column, 2)))
    data_start = min_col + 1 if spec.get('categories_col') else min_col
    data = Reference(ws, min_col=data_start, max_col=max_col, min_row=min_row, max_row=max_row)
    chart.add_data(data, titles_from_data=bool(spec.get('titles_from_data', True)))
    cat_col = int(spec.get('categories_col') or min_col)
    cats = Reference(ws, min_col=cat_col, min_row=min_row + (1 if spec.get('titles_from_data', True) else 0), max_row=max_row)
    chart.set_categories(cats)
    chart.height = float(spec.get('height', 8)); chart.width = float(spec.get('width', 14))
    ws.add_chart(chart, spec.get('anchor') or 'H2')


def _excel_apply_conditional(ws, spec):
    from openpyxl.formatting.rule import ColorScaleRule, CellIsRule
    from openpyxl.styles import PatternFill
    cell_range = spec.get('range')
    if not cell_range: return
    kind = str(spec.get('type') or 'color_scale').lower()
    if kind == 'cell':
        fill = PatternFill('solid', fgColor=_hex_color(spec.get('color'), 'FFC7CE'))
        ws.conditional_formatting.add(cell_range, CellIsRule(operator=spec.get('operator') or 'greaterThan', formula=[str(spec.get('value', 0))], fill=fill))
    else:
        ws.conditional_formatting.add(cell_range, ColorScaleRule(start_type='min', start_color='F8696B', mid_type='percentile', mid_value=50, mid_color='FFEB84', end_type='max', end_color='63BE7B'))


def _excel_fill_sheet(ws, spec, primary='1F4E78'):
    headers = spec.get('headers', []); rows = spec.get('rows', [])
    if headers:
        ws.append(headers); _excel_style_header(ws, 1, 1, len(headers), primary)
    for row in rows: ws.append(row)
    formulas = spec.get('formulas') or []
    for f in formulas:
        if f.get('cell'): ws[f['cell']] = f.get('formula')
    ws.freeze_panes = spec.get('freeze_panes', 'A2' if headers else None)
    if spec.get('autofilter', bool(headers)) and ws.max_row and ws.max_column: ws.auto_filter.ref = ws.dimensions
    for cf in spec.get('conditional_formats') or []: _excel_apply_conditional(ws, cf)
    for chart in spec.get('charts') or []: _excel_add_chart(ws, chart)
    if spec.get('rtl'): ws.sheet_view.rightToLeft = True
    if spec.get('autofit', True): _excel_autofit(ws)


def create_excel(payload):
    import openpyxl
    path = Path(payload['path']); path.parent.mkdir(parents=True, exist_ok=True)
    template = payload.get('template_path')
    wb = openpyxl.load_workbook(str(template), keep_vba=str(template).lower().endswith('.xlsm')) if template and Path(template).exists() else openpyxl.Workbook()
    if not template and wb.active: wb.remove(wb.active)
    primary = _hex_color((payload.get('branding') or {}).get('primary_color'), '1F4E78')
    for si, spec in enumerate(payload.get('sheets', [])):
        name = (spec.get('name') or f'Sheet{si+1}')[:31]
        if name in wb.sheetnames:
            ws = wb[name]
            if spec.get('replace_existing', True):
                for row in ws.iter_rows():
                    for cell in row: cell.value = None
        else: ws = wb.create_sheet(title=name)
        _excel_fill_sheet(ws, spec, primary)
    if not wb.sheetnames: wb.create_sheet('Sheet1')
    if payload.get('active_sheet') in wb.sheetnames: wb.active = wb.sheetnames.index(payload['active_sheet'])
    wb.save(str(path)); names = list(wb.sheetnames); wb.close()
    return {"success":True,"action":"create_excel","path":str(path),"sheets":names,"template_used":bool(template and Path(template).exists())}


def edit_excel(payload):
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment
    path=Path(payload['path']).resolve()
    if not path.exists(): raise FileNotFoundError(str(path))
    out_path=Path(payload.get('output_path') or path).resolve(); out_path.parent.mkdir(parents=True, exist_ok=True)
    keep_vba=path.suffix.lower()=='.xlsm'
    wb=openpyxl.load_workbook(str(path), keep_vba=keep_vba)
    primary=_hex_color((payload.get('branding') or {}).get('primary_color'), '1F4E78')
    applied=0
    for op in payload.get('operations') or []:
        typ=str(op.get('type') or '').lower(); sheet=op.get('sheet') or (wb.sheetnames[0] if wb.sheetnames else 'Sheet1')
        if typ=='add_sheet':
            if sheet not in wb.sheetnames: wb.create_sheet(sheet[:31]); applied+=1
            continue
        if sheet not in wb.sheetnames: wb.create_sheet(sheet[:31])
        ws=wb[sheet[:31]]
        if typ in ('set','set_value','set_formula'):
            cell=op.get('cell');
            if cell: ws[cell]=op.get('formula') if typ=='set_formula' else op.get('value'); applied+=1
        elif typ=='append_row': ws.append(op.get('values') or []); applied+=1
        elif typ=='rename_sheet': ws.title=str(op.get('new_name') or ws.title)[:31]; applied+=1
        elif typ=='freeze': ws.freeze_panes=op.get('cell') or 'A2'; applied+=1
        elif typ=='rtl': ws.sheet_view.rightToLeft=bool(op.get('value',True)); applied+=1
        elif typ=='autofilter': ws.auto_filter.ref=op.get('range') or ws.dimensions; applied+=1
        elif typ=='chart': _excel_add_chart(ws, op); applied+=1
        elif typ=='conditional_format': _excel_apply_conditional(ws, op); applied+=1
        elif typ=='format_range':
            for row in ws[op.get('range') or ws.dimensions]:
                for c in row:
                    if op.get('bold') is not None: c.font=Font(name=c.font.name, size=c.font.sz, bold=bool(op.get('bold')), italic=c.font.italic, color=c.font.color)
                    if op.get('fill'): c.fill=PatternFill('solid', fgColor=_hex_color(op.get('fill')))
                    if op.get('align'): c.alignment=Alignment(horizontal=op.get('align'), vertical='center')
            applied+=1
        elif typ=='autofit': _excel_autofit(ws); applied+=1
    if payload.get('recalculate_on_open', True):
        try:
            wb.calculation.fullCalcOnLoad=True; wb.calculation.forceFullCalc=True; wb.calculation.calcMode='auto'
        except Exception: pass
    wb.save(str(out_path)); names=list(wb.sheetnames); wb.close()
    return {"success":True,"action":"edit_excel","path":str(out_path),"operations_applied":applied,"sheets":names}


def _ppt_add_bullets(slide, bullets, font_size=22):
    from pptx.util import Pt
    placeholder = None
    for sh in slide.placeholders:
        if getattr(sh, 'has_text_frame', False) and sh != slide.shapes.title:
            placeholder = sh; break
    if placeholder is None: return
    tf=placeholder.text_frame; tf.clear()
    for idx, bullet in enumerate(bullets or []):
        p=tf.paragraphs[0] if idx==0 else tf.add_paragraph(); p.text=safe_text(bullet); p.level=0
        for run in p.runs: run.font.size=Pt(font_size)


def _ppt_add_table(slide, spec):
    from pptx.util import Inches
    headers=spec.get('headers') or []; rows=spec.get('rows') or []
    cols=max(len(headers), max((len(r) for r in rows), default=0),1); total_rows=len(rows)+(1 if headers else 0)
    shape=slide.shapes.add_table(max(total_rows,1), cols, Inches(float(spec.get('left',0.7))), Inches(float(spec.get('top',2.0))), Inches(float(spec.get('width',11.9))), Inches(float(spec.get('height',3.7))))
    table=shape.table; offset=0
    if headers:
        for c,v in enumerate(headers[:cols]): table.cell(0,c).text=safe_text(v)
        offset=1
    for r,row in enumerate(rows):
        for c in range(cols): table.cell(r+offset,c).text=safe_text(row[c] if c<len(row) else '')
    return table


def create_pptx(payload):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    path=Path(payload['path']); path.parent.mkdir(parents=True, exist_ok=True)
    template=payload.get('template_path')
    prs=Presentation(str(template)) if template and Path(template).exists() else Presentation()
    if not template:
        prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    brand=payload.get('branding') or {}; primary=_hex_color(brand.get('primary_color'),'1F4E78')
    if payload.get('include_title_slide', True):
        slide=prs.slides.add_slide(prs.slide_layouts[0]); slide.shapes.title.text=payload.get('title','Presentation')
        if len(slide.placeholders)>1: slide.placeholders[1].text=payload.get('subtitle') or 'Generated by ABDULKAREM AI X'
    for spec in payload.get('slides',[]):
        layout_index=int(spec.get('layout',1)); layout_index=min(max(layout_index,0),len(prs.slide_layouts)-1)
        slide=prs.slides.add_slide(prs.slide_layouts[layout_index])
        if slide.shapes.title: slide.shapes.title.text=spec.get('title','')
        _ppt_add_bullets(slide,spec.get('bullets',[]),int(spec.get('font_size',22)))
        for tb in spec.get('tables') or []: _ppt_add_table(slide,tb)
        for im in spec.get('images') or []:
            image_path=im.get('path') if isinstance(im,dict) else im
            if image_path and Path(image_path).exists():
                slide.shapes.add_picture(str(image_path), Inches(float(im.get('left',7.4) if isinstance(im,dict) else 7.4)), Inches(float(im.get('top',1.7) if isinstance(im,dict) else 1.7)), width=Inches(float(im.get('width',5.2) if isinstance(im,dict) else 5.2)))
        notes=spec.get('notes')
        if notes:
            try: slide.notes_slide.notes_text_frame.text=safe_text(notes)
            except Exception: pass
    prs.save(str(path)); count=len(prs.slides)
    return {"success":True,"action":"create_pptx","path":str(path),"slides":count,"template_used":bool(template and Path(template).exists())}


def edit_pptx(payload):
    from pptx import Presentation
    from pptx.util import Inches
    path=Path(payload['path']).resolve()
    if not path.exists(): raise FileNotFoundError(str(path))
    out_path=Path(payload.get('output_path') or path).resolve(); out_path.parent.mkdir(parents=True,exist_ok=True)
    prs=Presentation(str(path)); applied=0
    for op in payload.get('operations') or []:
        typ=str(op.get('type') or '').lower()
        if typ=='add_slide':
            idx=min(max(int(op.get('layout',1)),0),len(prs.slide_layouts)-1); slide=prs.slides.add_slide(prs.slide_layouts[idx])
            if slide.shapes.title: slide.shapes.title.text=op.get('title','')
            _ppt_add_bullets(slide,op.get('bullets',[]),int(op.get('font_size',22)))
            for tb in op.get('tables') or []: _ppt_add_table(slide,tb)
            applied+=1; continue
        slide_no=int(op.get('slide',1));
        if slide_no<1 or slide_no>len(prs.slides): continue
        slide=prs.slides[slide_no-1]
        if typ=='set_title' and slide.shapes.title: slide.shapes.title.text=safe_text(op.get('title')); applied+=1
        elif typ=='replace_text':
            old,new=safe_text(op.get('find')),safe_text(op.get('replace'))
            for sh in slide.shapes:
                if getattr(sh,'has_text_frame',False):
                    for p in sh.text_frame.paragraphs:
                        for run in p.runs:
                            if old and old in run.text: run.text=run.text.replace(old,new); applied+=1
        elif typ=='add_image':
            ip=op.get('path')
            if ip and Path(ip).exists(): slide.shapes.add_picture(str(ip), Inches(float(op.get('left',7.4))), Inches(float(op.get('top',1.7))), width=Inches(float(op.get('width',5.2)))); applied+=1
        elif typ=='add_table': _ppt_add_table(slide,op); applied+=1
        elif typ=='set_notes':
            try: slide.notes_slide.notes_text_frame.text=safe_text(op.get('notes')); applied+=1
            except Exception: pass
    prs.save(str(out_path)); count=len(prs.slides)
    return {"success":True,"action":"edit_pptx","path":str(out_path),"operations_applied":applied,"slides":count}

def export_pdf(payload):
    if os.name != 'nt':
        raise RuntimeError('Office COM PDF export يعمل على Windows فقط مع Microsoft Office مثبت.')
    try:
        import win32com.client
    except Exception as e:
        raise RuntimeError('pywin32 غير مثبت. نفّذ: pip install pywin32') from e
    path = Path(payload['path']).resolve()
    if not path.exists(): raise FileNotFoundError(str(path))
    pdf = path.with_suffix('.pdf'); ext = path.suffix.lower(); app = None
    try:
        if ext in ('.doc','.docx'):
            app = win32com.client.DispatchEx('Word.Application'); app.Visible = False
            doc = app.Documents.Open(str(path)); doc.ExportAsFixedFormat(str(pdf), 17); doc.Close(False)
        elif ext in ('.ppt','.pptx'):
            app = win32com.client.DispatchEx('PowerPoint.Application')
            prs = app.Presentations.Open(str(path), WithWindow=False); prs.SaveAs(str(pdf), 32); prs.Close()
        elif ext in ('.xls','.xlsx','.xlsm'):
            app = win32com.client.DispatchEx('Excel.Application'); app.Visible = False
            wb = app.Workbooks.Open(str(path)); wb.RefreshAll(); app.CalculateFull(); wb.ExportAsFixedFormat(0, str(pdf)); wb.Close(False)
        else:
            raise RuntimeError('التحويل إلى PDF يدعم Word/Excel/PowerPoint فقط.')
    finally:
        try:
            if app: app.Quit()
        except Exception: pass
    return {"success":True,"action":"export_pdf","path":str(pdf)}


def main():
    payload = json.loads(sys.stdin.read() or '{}')
    action = payload.get('action')
    if action == 'inspect': result = inspect_file(payload['path'])
    elif action == 'compare': result = compare_files(payload)
    elif action == 'search': result = search_files(payload)
    elif action == 'render_pdf_pages': result = render_pdf_pages(payload)
    elif action == 'create_word': result = create_word(payload)
    elif action == 'create_excel': result = create_excel(payload)
    elif action == 'create_pptx': result = create_pptx(payload)
    elif action == 'edit_word': result = edit_word(payload)
    elif action == 'edit_excel': result = edit_excel(payload)
    elif action == 'edit_pptx': result = edit_pptx(payload)
    elif action == 'export_pdf': result = export_pdf(payload)
    else: raise RuntimeError(f'Unknown action: {action}')
    out(result)


if __name__ == '__main__':
    try:
        main()
    except Exception as e:
        sys.stderr.write(f"{type(e).__name__}: {e}\n")
        sys.stderr.write(traceback.format_exc())
        sys.exit(1)
